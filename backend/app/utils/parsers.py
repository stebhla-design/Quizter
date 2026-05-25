import re
import urllib.parse
from typing import Union
import requests
from bs4 import BeautifulSoup

def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF file bytes."""
    try:
        import io
        from pypdf import PdfReader
        
        pdf_file = io.BytesIO(file_content)
        reader = PdfReader(pdf_file)
        text = ""
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error parsing PDF: {e}")
        return ""

def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from Word .docx file bytes."""
    try:
        import io
        import docx
        
        doc_file = io.BytesIO(file_content)
        doc = docx.Document(doc_file)
        text = ""
        for paragraph in doc.paragraphs:
            if paragraph.text:
                text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error parsing DOCX: {e}")
        return ""

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PowerPoint .pptx file bytes."""
    try:
        import io
        from pptx import Presentation
        
        pptx_file = io.BytesIO(file_content)
        prs = Presentation(pptx_file)
        text = ""
        for slide in prs.slides:
            for shape in prs.slides[slide.slide_id - prs.slides[0].slide_id if prs.slides else 0].shapes if hasattr(prs, 'slides') else []:
                pass
            
            # Use slide shapes safely
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error parsing PPTX: {e}")
        return ""

def extract_text_from_url(url: str) -> str:
    """Fetch URL and parse HTML contents to return purified text."""
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }
        
        if not url.startswith(("http://", "https://")):
            url = "https://" + url

        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        soup = BeautifulSoup(response.content, "html.parser")
        
        # Strip script, style, head, nav, footer
        for element in soup(["script", "style", "nav", "footer", "header", "aside"]):
            element.decompose()
            
        title = soup.title.string.strip() if soup.title and soup.title.string else ""
        body_text = soup.get_text(separator="\n")
        
        # Clean extra spacing
        lines = [line.strip() for line in body_text.splitlines() if line.strip()]
        cleaned_text = "\n".join(lines)
        
        if title:
            return f"Title: {title}\n\nContent:\n{cleaned_text}"
        return cleaned_text
    except Exception as e:
        print(f"Error fetching/parsing URL {url}: {e}")
        return ""
